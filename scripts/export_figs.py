from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("export_figs")


FIG_DIR = Path("data/figs")
SLIDES_OUT = Path("slides/als-perturb-slides.pptx")


def _add_title(prs, title: str):
    slide_layout = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    return slide


def _add_bullets(txBox, bullets):
    tf = txBox.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
        p.text = b
        p.level = 0


def build_deck():
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    SLIDES_OUT.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    # Slide 1: Task 1 workflow
    s1 = _add_title(prs, "Task 1: In-Silico Perturbation Workflow")
    left = Inches(0.5); top = Inches(1.5); width = Inches(9)
    txBox = s1.shapes.add_textbox(left, top, width, Inches(1.5))
    _add_bullets(txBox, [
        "Rank-based swaps and expression-based log2 fold-change",
        "Safety checks: bounds, renormalization, missing genes skipped",
        "Outputs: perturbed AnnData, figures under data/figs/",
    ])
    wf = FIG_DIR / "task1_workflow.png"
    if wf.exists():
        s1.shapes.add_picture(str(wf), Inches(0.5), Inches(3), height=Inches(3))

    # Slide 2: Task 2 embedding snapshot
    s2 = _add_title(prs, "Task 2: Embedding Snapshot (Geneformer V2)")
    fig2 = FIG_DIR / "task2_embedding_snapshot.png"
    if fig2.exists():
        s2.shapes.add_picture(str(fig2), Inches(0.5), Inches(1.5), height=Inches(5))
    else:
        tx = s2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
        _add_bullets(tx, ["Run notebook 02 to generate snapshot figure."])

    # Slide 3: Task 3 UMAP + metrics
    s3 = _add_title(prs, "Task 3: Interpretation and Metrics")
    um = FIG_DIR / "task3_umap.png"
    cs = FIG_DIR / "task3_centroid_shifts.png"
    if um.exists():
        s3.shapes.add_picture(str(um), Inches(0.5), Inches(1.5), height=Inches(3))
    if cs.exists():
        s3.shapes.add_picture(str(cs), Inches(5.0), Inches(1.5), height=Inches(3))
    mt = FIG_DIR / "task3_metrics.csv"
    if mt.exists():
        tx = s3.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2))
        _add_bullets(tx, [f"Metrics table saved at {mt}"])

    # Slide 4: Task 4 Top targets (optional)
    s4 = _add_title(prs, "Task 4 (Optional): Top Targets")
    fig4 = FIG_DIR / "task4_top_targets.png"
    if fig4.exists():
        s4.shapes.add_picture(str(fig4), Inches(0.5), Inches(1.5), height=Inches(5))
    else:
        tx = s4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
        _add_bullets(tx, ["Run notebook 04 to generate target ranking figure."])

    prs.save(str(SLIDES_OUT))
    logger.info("Saved slides to %s", SLIDES_OUT)


if __name__ == "__main__":
    build_deck()


